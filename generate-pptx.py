"""
Generate Providence APIM & AI Gateway Deep Dive Presentation
=============================================================
Creates a professional .pptx with official Microsoft Learn diagrams and
polished visuals. Images downloaded from:
  https://learn.microsoft.com/en-us/azure/api-management/genai-gateway-capabilities

Covers:
  - AI Gateway capabilities overview (MS Learn diagram)
  - APIM vs AI Gateway education
  - Token rate limiting (MS Learn diagram)
  - Semantic caching (MS Learn diagram)
  - Load balancing & circuit breaker (MS Learn diagrams)
  - Content safety (MS Learn diagram)
  - Token metrics / observability (MS Learn diagram)
  - Content Understanding patient discharge use case
  - Budget & cost controls
  - CI/CD pipeline
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os


# ─── Paths ───────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(SCRIPT_DIR, "images")

# ─── Color Palette (Microsoft Azure) ────────────────────
DARK_BG    = RGBColor(0x0F, 0x11, 0x16)   # near-black
BLUE       = RGBColor(0x00, 0x78, 0xD4)   # Azure blue
LIGHT_BLUE = RGBColor(0x50, 0xE6, 0xFF)
ACCENT_BLUE = RGBColor(0x00, 0x50, 0x8F)  # darker accent
GREEN      = RGBColor(0x10, 0x7C, 0x10)
TEAL       = RGBColor(0x00, 0xB7, 0xC3)
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)
RED        = RGBColor(0xD1, 0x34, 0x38)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
BG_GRAY    = RGBColor(0xF5, 0xF5, 0xF5)
CARD_BG    = RGBColor(0xF0, 0xF6, 0xFC)   # light blue card
BORDER_BLUE = RGBColor(0xA0, 0xC4, 0xE8)


def img(name):
    """Return full path for an image in the images/ folder."""
    return os.path.join(IMG_DIR, name)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bg(slide, color_top, color_bottom):
    """Set a two-stop gradient background."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color_top
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color_bottom
    fill.gradient_stops[1].position = 1.0


def add_shape_with_text(slide, left, top, width, height, text, font_size=14,
                        bold=False, color=WHITE, bg_color=None, alignment=PP_ALIGN.LEFT,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Add a shape with text."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return shape


def add_card(slide, left, top, width, height, title, body_lines,
             title_bg=BLUE, body_bg=CARD_BG, title_size=16, body_size=13,
             title_color=WHITE, body_color=DARK_GRAY):
    """Add a card with colored header and body with bullet lines."""
    add_shape_with_text(slide, left, top, width, Inches(0.5),
                        title, font_size=title_size, bold=True,
                        color=title_color, bg_color=title_bg,
                        alignment=PP_ALIGN.CENTER)
    body_shape = add_shape_with_text(
        slide, left, top + Inches(0.5), width, height - Inches(0.5),
        "", font_size=body_size, color=body_color, bg_color=body_bg)
    tf = body_shape.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(body_size)
        p.font.color.rgb = body_color
        p.space_after = Pt(4)
    return body_shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_list(text_frame, items, font_size=13, color=DARK_GRAY, bold_first=False):
    """Add bullet list items to a text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
        p.level = 0


def add_connector_arrow(slide, x1, y1, x2, y2, color=BLUE, width=Pt(3)):
    """Add a line connector between two points."""
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    connector.line.fill.solid()
    connector.line.fill.fore_color.rgb = color


def add_section_divider(slide, left, top, width, text, color=BLUE):
    """Add a thin accent bar with section text."""
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text_box(slide, left + Inches(0.25), top, width, Inches(0.5),
                 text, font_size=20, color=color, bold=True)


def set_notes(slide, notes_text):
    """Set speaker notes for a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ══════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(1.4), Inches(0.08), Inches(3.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    add_text_box(slide, Inches(1.0), Inches(1.4), Inches(7), Inches(0.7),
                 "Azure API Management", font_size=20, color=LIGHT_BLUE, bold=True)
    add_text_box(slide, Inches(1.0), Inches(2.2), Inches(7), Inches(1.5),
                 "AI Gateway Deep Dive", font_size=48, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.0), Inches(3.8), Inches(7), Inches(0.8),
                 "Govern  ·  Observe  ·  Optimize  Your AI Workloads",
                 font_size=22, color=LIGHT_GRAY)
    add_text_box(slide, Inches(1.0), Inches(5.3), Inches(7), Inches(0.4),
                 "Providence Health & Services  |  Microsoft Azure  |  April 2026",
                 font_size=15, color=MED_GRAY)

    # Right side — capabilities overview image
    try:
        slide.shapes.add_picture(img("capabilities-summary.png"),
                                  Inches(8.3), Inches(1.2), Inches(4.5))
    except Exception:
        pass

    set_notes(slide, """SPEAKER NOTES — Title Slide
Welcome everyone. Today we'll do a deep dive into Azure API Management as an AI Gateway.
We'll go from zero governance to a fully managed, observable, cost-controlled AI layer.
The diagram on the right is from official Microsoft documentation showing all AI Gateway capabilities.
By the end, you'll see: direct vs governed calls, semantic caching savings, Content Understanding for patient discharge docs, and a CI/CD pipeline.
Duration: ~90 minutes with live demos.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 2 — Agenda
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title with accent bar
    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10), "Session Agenda", DARK_BG)

    agenda_items = [
        ("1", "The Problem — Direct Calls", "15 min", "Zero governance, API key sprawl, no cost visibility", RED),
        ("2", "Education — APIM as AI Gateway", "20 min", "AI Gateway capabilities, MS Learn architecture diagrams", BLUE),
        ("3", "Hands-On — Setup the Gateway", "20 min", "Provision APIM, backends, AI-aware policies", TEAL),
        ("4", "The Transformation — Governed Calls", "15 min", "Token metrics, rate limiting, semantic caching", GREEN),
        ("5", "Content Understanding", "10 min", "Patient discharge docs through the same gateway", ORANGE),
        ("6", "Budget, CI/CD & Next Steps", "10 min", "$25 budget alerts, GitHub Actions pipeline", ACCENT_BLUE),
    ]

    y_start = Inches(1.3)
    for i, (num, title, duration, desc, clr) in enumerate(agenda_items):
        y = y_start + Inches(i * 0.95)
        # Number circle
        add_shape_with_text(slide, Inches(0.8), y, Inches(0.6), Inches(0.6),
                            num, font_size=22, bold=True, color=WHITE, bg_color=clr,
                            alignment=PP_ALIGN.CENTER, shape_type=MSO_SHAPE.OVAL)
        # Title
        add_text_box(slide, Inches(1.6), y + Inches(0.02), Inches(6.5), Inches(0.35),
                     title, font_size=18, bold=True, color=DARK_BG)
        # Description
        add_text_box(slide, Inches(1.6), y + Inches(0.38), Inches(6.5), Inches(0.3),
                     desc, font_size=12, color=MED_GRAY)
        # Duration pill
        add_shape_with_text(slide, Inches(9.0), y + Inches(0.08), Inches(1.3), Inches(0.4),
                            duration, font_size=13, bold=True, color=clr, bg_color=BG_GRAY,
                            alignment=PP_ALIGN.CENTER)

    # Bottom image strip — capabilities overview
    try:
        slide.shapes.add_picture(img("capabilities-summary.png"),
                                  Inches(10.5), Inches(1.5), Inches(2.5))
    except Exception:
        pass

    add_text_box(slide, Inches(10.5), Inches(0.7), Inches(2.5), Inches(0.6),
                 "AI Gateway\nCapabilities", font_size=12, color=BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Agenda
Six parts today. We start with the problem — direct calls with no guardrails.
Then education on what APIM is and how AI Gateway extends it — using official Microsoft Learn diagrams.
Live hands-on setup, then see the transformation in action.
NEW today: Content Understanding for patient discharge documents — showing APIM governs ALL AI services, not just LLMs.
We close with budget alerts at $25/month and a CI/CD pipeline on GitHub.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 3 — The Problem (No APIM)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Part 1: The Problem — Direct Azure OpenAI Calls", RED)

    # LEFT — visual architecture: 3 app boxes with arrows to OpenAI box
    app_names = ["Patient Portal", "Clinical Tool", "Admin Dashboard"]
    app_colors = [RGBColor(0x8B, 0x00, 0x00), RGBColor(0x8B, 0x00, 0x00), RGBColor(0x8B, 0x00, 0x00)]
    for i, name in enumerate(app_names):
        y = Inches(1.5 + i * 1.15)
        # App box
        add_shape_with_text(slide, Inches(0.5), y, Inches(2.2), Inches(0.85),
                            f"🖥️  {name}\nAPI Key: sk-xxx...{i+1}",
                            font_size=12, bold=False, color=WHITE, bg_color=RED,
                            alignment=PP_ALIGN.CENTER)
        # Arrow text
        add_text_box(slide, Inches(2.75), y + Inches(0.2), Inches(1.5), Inches(0.4),
                     "── API Key ──▶", font_size=11, color=RED, bold=True)

    # OpenAI box
    add_shape_with_text(slide, Inches(4.3), Inches(1.8), Inches(2.0), Inches(2.5),
                        "Azure\nOpenAI\n(gpt-4o)\n\nNo rate limits\nNo caching\nNo metrics",
                        font_size=12, bold=True, color=WHITE, bg_color=DARK_GRAY,
                        alignment=PP_ALIGN.CENTER)

    # RIGHT — risk cards
    risk_data = [
        ("🔑  Key Sprawl", "Every app stores API keys.\nNo rotation, no revocation."),
        ("💸  Cost Blindness", "No per-app attribution.\nDuplicate queries = duplicate cost."),
        ("🚫  No Rate Limits", "One runaway app burns\nall tokens for everyone."),
        ("⚠️  No Safety", "No content filtering.\nNo audit trail."),
    ]
    for i, (title, desc) in enumerate(risk_data):
        x = Inches(6.8) + Inches((i % 2) * 3.2)
        y = Inches(1.3) + Inches((i // 2) * 2.3)
        add_card(slide, x, y, Inches(3.0), Inches(2.0),
                 title, desc.split("\n"),
                 title_bg=RED, body_bg=RGBColor(0xFD, 0xED, 0xED),
                 body_color=DARK_GRAY)

    # Demo callout
    add_shape_with_text(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8),
                        "🔴  LIVE DEMO: 01-without-apim.py — 10 rapid requests, zero governance, watch the tokens burn",
                        font_size=16, bold=True, color=WHITE, bg_color=RED,
                        alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — The Problem
This is how most teams start. Direct calls to Azure OpenAI with API keys embedded in each app.
DEMO: Run 01-without-apim.py. Show 10 rapid requests — no rate limiting, no caching.
Point out: total tokens consumed, no observability, same question 5 times = 5x the cost.
Each app has its own API key — no centralized revocation, no rotation policy.
Key message: 'This is the burning platform. Let's fix it.'""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 4 — What is APIM? (with capabilities overview image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Part 2: What is Azure API Management?", BLUE)

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
                 "A fully managed service for publishing, securing, transforming, and monitoring APIs",
                 font_size=17, color=MED_GRAY)

    # Three pillar cards
    pillars = [
        ("API Gateway", ["Route & transform requests", "Rate-limit & throttle", "Authenticate callers"], BLUE),
        ("Developer Portal", ["Self-service API docs", "Interactive try-it console", "Onboard new teams"], GREEN),
        ("Management Plane", ["Policy configuration", "Analytics dashboards", "Lifecycle management"], ORANGE),
    ]
    for i, (title, items, clr) in enumerate(pillars):
        x = Inches(0.5 + i * 4.2)
        add_card(slide, x, Inches(1.8), Inches(3.9), Inches(2.3),
                 title, items, title_bg=clr, body_bg=CARD_BG)

    # KEY INSIGHT — big box with image
    add_shape_with_text(slide, Inches(0.5), Inches(4.4), Inches(7.5), Inches(2.6),
                        "KEY INSIGHT\n\n"
                        "APIM = General API Gateway  (any REST / SOAP / GraphQL API)\n"
                        "AI Gateway = APIM + AI-Aware Policies\n"
                        "    • Token-based rate limiting (not just request-count)\n"
                        "    • Semantic caching (not just exact URL match)\n"
                        "    • LLM content safety filters\n"
                        "    • Token consumption metrics\n\n"
                        "   It's NOT a separate product — it's APIM with intelligence about LLM semantics.",
                        font_size=14, color=WHITE, bg_color=BLUE,
                        alignment=PP_ALIGN.LEFT)

    # Foundry import screenshot
    try:
        slide.shapes.add_picture(img("ai-foundry-import.png"),
                                  Inches(8.3), Inches(4.4), Inches(4.5))
    except Exception:
        pass
    add_text_box(slide, Inches(8.3), Inches(7.0), Inches(4.5), Inches(0.3),
                 "Foundry model import in Azure Portal", font_size=10, color=MED_GRAY,
                 alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — What is APIM
Azure API Management has three components: Gateway, Developer Portal, Management Plane.
The key insight: AI Gateway is NOT a separate product. It's APIM configured with AI-specific policies.
Standard APIM counts requests per minute. AI Gateway counts TOKENS per minute — that's the real cost driver.
Standard APIM caches by exact URL. AI Gateway uses semantic similarity.
Screenshot shows the Foundry model import wizard — streamlined onboarding of LLM endpoints.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 5 — AI Gateway Capabilities Overview (full MS Learn image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "AI Gateway Capabilities — Official Overview", BLUE)

    # Full-width capabilities image
    try:
        slide.shapes.add_picture(img("capabilities-summary.png"),
                                  Inches(0.8), Inches(1.2), Inches(11.5))
    except Exception:
        pass

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
                 "Source: learn.microsoft.com/azure/api-management/genai-gateway-capabilities",
                 font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Capabilities Overview
This is the official Microsoft Learn diagram showing ALL AI Gateway capabilities.
Walk through each section:
- Traffic mediation: Import models from Foundry, manage MCP servers, A2A agents
- Scalability: Token rate limiting, semantic caching, native scaling
- Security: Managed identity auth, OAuth, content safety
- Resiliency: Load balancing (round-robin, weighted, priority), circuit breaker
- Observability: Token metrics, prompt/completion logging, dashboards
- Developer Experience: API Center integration, self-service portals
This is what we are implementing today.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 6 — Comparison: Standard APIM vs AI Gateway
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Standard APIM  vs.  APIM as AI Gateway", DARK_BG)

    # Table header
    col_x = [Inches(0.8), Inches(4.2), Inches(7.3), Inches(10.5)]
    col_w = [Inches(3.2), Inches(2.9), Inches(3.0), Inches(2.3)]
    headers = ["Capability", "Standard APIM", "AI Gateway", "Impact"]
    for x, w, h in zip(col_x, col_w, headers):
        add_shape_with_text(slide, x, Inches(1.2), w, Inches(0.5),
                            h, font_size=14, bold=True, color=WHITE, bg_color=DARK_BG,
                            alignment=PP_ALIGN.CENTER)

    rows = [
        ("Rate Limiting", "Requests / minute", "Tokens / minute (TPM)", "💰 Fair cost split"),
        ("Caching", "Exact URL match", "Semantic similarity", "💰 60-80% savings"),
        ("Metrics", "Request count, latency", "Token consumption / model", "📊 Cost attribution"),
        ("Content Filtering", "None built-in", "LLM content safety", "🛡️ Compliance"),
        ("Cost Control", "Basic rate limits", "Token budgets / consumer", "💰 Predictable spend"),
        ("Load Balancing", "Round-robin / weighted", "429-aware + circuit breaker", "🔄 Zero downtime"),
        ("Authentication", "Keys, OAuth, certs", "+ Managed Identity → AI", "🔐 No key sprawl"),
    ]

    for i, (cap, std, ai, impact) in enumerate(rows):
        y = Inches(1.85 + i * 0.6)
        bg = BG_GRAY if i % 2 == 0 else WHITE
        vals = [cap, std, ai, impact]
        bolds = [True, False, True, False]
        colors = [DARK_BG, MED_GRAY, GREEN, DARK_GRAY]
        for x, w, val, b, c in zip(col_x, col_w, vals, bolds, colors):
            add_shape_with_text(slide, x, y, w, Inches(0.5),
                                val, font_size=12, bold=b, color=c, bg_color=bg,
                                alignment=PP_ALIGN.CENTER)

    # Bottom callout
    add_shape_with_text(slide, Inches(0.8), Inches(6.3), Inches(12), Inches(0.8),
                        "The AI Gateway transforms APIM from a request counter to a token-aware, semantically intelligent AI cost manager",
                        font_size=15, bold=True, color=WHITE, bg_color=BLUE,
                        alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Comparison Table
Walk through each row. The key differentiators:
1. Rate limiting by TOKENS not requests — because a 1-token request and a 10K-token request cost very differently.
2. Semantic caching — biggest cost saver. 'Symptoms of pneumonia' matches 'signs of pneumonia'.
3. Token metrics — emit per-subscription, per-model consumption to Azure Monitor.
4. Content safety — blocks jailbreak attempts, hate speech BEFORE the model.
5. 429-aware load balancing — when one backend is throttled, route to another automatically.
6. Managed identity — APIM authenticates to Azure OpenAI. No API keys stored anywhere.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 7 — Architecture With APIM AI Gateway
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Architecture: With APIM AI Gateway", BLUE)

    # LEFT — Apps column
    apps = [
        ("Patient Portal", "Subscription Key A"),
        ("Clinical Tool", "Subscription Key B"),
        ("Admin Dashboard", "Subscription Key C"),
    ]
    for i, (name, key) in enumerate(apps):
        y = Inches(1.8 + i * 1.4)
        add_shape_with_text(slide, Inches(0.3), y, Inches(2.3), Inches(1.0),
                            f"🖥️  {name}\n{key}",
                            font_size=12, bold=False, color=WHITE, bg_color=ACCENT_BLUE,
                            alignment=PP_ALIGN.CENTER)
        # Arrow
        add_text_box(slide, Inches(2.65), y + Inches(0.25), Inches(0.6), Inches(0.4),
                     "▶", font_size=20, color=BLUE, bold=True)

    # CENTER — APIM Gateway (large blue box)
    gw_text = ("APIM AI Gateway\n"
               "━━━━━━━━━━━━━━━━\n"
               "🔐  Managed Identity Auth\n"
               "📊  Token Metrics Emission\n"
               "⏱️  Token Rate Limiting\n"
               "💾  Semantic Cache Lookup\n"
               "🛡️  Content Safety Filter\n"
               "🔄  429-Aware Load Balancer")
    add_shape_with_text(slide, Inches(3.3), Inches(1.3), Inches(4.0), Inches(5.0),
                        gw_text, font_size=14, color=WHITE, bg_color=BLUE,
                        alignment=PP_ALIGN.CENTER)

    # RIGHT — Backends
    backends = [
        ("Azure OpenAI\ngpt-4o  •  eastus2", GREEN),
        ("Content\nUnderstanding\nwestus", ORANGE),
        ("Future: Speech\nVision, etc.", MED_GRAY),
    ]
    for i, (name, clr) in enumerate(backends):
        y = Inches(1.5 + i * 1.8)
        # Arrow from gateway
        add_text_box(slide, Inches(7.35), y + Inches(0.3), Inches(0.6), Inches(0.4),
                     "▶", font_size=20, color=clr, bold=True)
        add_shape_with_text(slide, Inches(7.9), y, Inches(2.8), Inches(1.2),
                            name, font_size=13, bold=True, color=WHITE, bg_color=clr,
                            alignment=PP_ALIGN.CENTER)

    # Bottom: load balancing diagram from MS Learn
    try:
        slide.shapes.add_picture(img("backend-load-balancing.png"),
                                  Inches(10.8), Inches(1.3), Inches(2.2))
    except Exception:
        pass
    add_text_box(slide, Inches(10.8), Inches(4.0), Inches(2.2), Inches(0.3),
                 "Backend Load Balancing", font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # Key callout
    add_shape_with_text(slide, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.6),
                        "One Gateway  →  One Subscription Key  →  All AI Services  →  One Dashboard",
                        font_size=16, bold=True, color=WHITE, bg_color=ACCENT_BLUE,
                        alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Architecture
Single entry point for ALL AI services. Apps use subscription keys — never direct AI API keys.
The gateway handles: authentication via managed identity, semantic cache lookup, token rate limiting, metrics, content safety.
KEY POINT: Multiple backends shown — Azure OpenAI AND Content Understanding. Same gateway, same policies, same dashboard.
The MS Learn diagram on the right shows how backend load balancing distributes across multiple deployments.
This is the power: one pane of glass for all AI consumption.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 8 — Token Rate Limiting (MS Learn image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "AI Policy: Token Rate Limiting", BLUE)

    # LEFT — explanation
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.4),
                 "Why tokens, not requests?", font_size=18, color=DARK_BG, bold=True)

    limit_items = [
        "• A 10-word prompt costs ~15 tokens",
        "• A 500-word prompt costs ~750 tokens",
        "• Same request count, 50× different cost!",
        "",
        "The llm-token-limit policy enforces TPM",
        "(Tokens Per Minute) per consumer key.",
        "",
        "• Pre-calculates prompt tokens gateway-side",
        "• Rejects over-limit before reaching backend",
        "• Returns 429 with Retry-After header",
    ]
    txBox = add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5), "")
    add_bullet_list(txBox.text_frame, limit_items, font_size=14, color=DARK_GRAY)

    # RIGHT — MS Learn diagram
    try:
        slide.shapes.add_picture(img("token-rate-limiting.png"),
                                  Inches(6.8), Inches(1.2), Inches(6.0))
    except Exception:
        pass
    add_text_box(slide, Inches(6.8), Inches(5.0), Inches(6.0), Inches(0.3),
                 "Source: Microsoft Learn — Token rate limiting in API Management",
                 font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # Policy XML snippet
    policy_xml = ('<llm-token-limit\n'
                  '    counter-key="@(context.Subscription.Id)"\n'
                  '    tokens-per-minute="10000"\n'
                  '    estimate-prompt-tokens="true"\n'
                  '    remaining-tokens-variable-name="remainingTokens" />')
    add_shape_with_text(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(1.5),
                        policy_xml, font_size=13, color=LIGHT_BLUE, bg_color=DARK_BG,
                        alignment=PP_ALIGN.LEFT)

    set_notes(slide, """SPEAKER NOTES — Token Rate Limiting
The MS Learn diagram on the right shows the token rate limiting flow.
Unlike request limits (100 requests/min), this counts TOKENS.  A single long prompt can consume thousands of tokens.
The policy can pre-calculate prompt tokens on the gateway side — so if you're already over limit, it rejects immediately without even calling the backend.
10K TPM per subscription gives predictable cost control per app/team.
Example: Patient Portal gets 10K TPM, Clinical Tool gets 20K TPM, Admin Dashboard gets 5K TPM.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 9 — Semantic Caching (MS Learn image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "AI Policy: Semantic Caching — Save 60-80%", GREEN)

    # MS Learn caching diagram — full width
    try:
        slide.shapes.add_picture(img("semantic-caching.png"),
                                  Inches(0.8), Inches(1.2), Inches(6.5))
    except Exception:
        pass
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(6.5), Inches(0.3),
                 "Source: Microsoft Learn — Semantic caching in API Management",
                 font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # RIGHT — how it works + examples
    add_text_box(slide, Inches(7.8), Inches(1.2), Inches(5.0), Inches(0.4),
                 "How It Works", font_size=18, color=GREEN, bold=True)
    flow_items = [
        "1. Prompt → Generate embedding vector",
        "2. Search Redis for similar cached prompts",
        "3. Similarity > 80%? → Return cached completion",
        "4. Cache miss → Call Azure OpenAI → Store result",
    ]
    txBox = add_text_box(slide, Inches(7.8), Inches(1.8), Inches(5.0), Inches(1.8), "")
    add_bullet_list(txBox.text_frame, flow_items, font_size=13, color=DARK_GRAY)

    add_text_box(slide, Inches(7.8), Inches(3.5), Inches(5.0), Inches(0.4),
                 "Same Intent, Different Words", font_size=16, color=GREEN, bold=True)
    examples = [
        '"What are pneumonia symptoms?"  → MISS (1st call)',
        '"What symptoms does pneumonia?"  → HIT  ✅',
        '"Signs of pneumonia"                    → HIT  ✅',
        '"How do I know if I have pneumonia?"  → HIT  ✅',
    ]
    txBox = add_text_box(slide, Inches(7.8), Inches(4.0), Inches(5.0), Inches(1.8), "")
    add_bullet_list(txBox.text_frame, examples, font_size=12, color=DARK_GRAY)

    # Bottom savings callout
    add_shape_with_text(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8),
                        "Patient FAQ (80% cache hit) → $360/mo savings  |  Clinical Decision Support (60% hit) → $108/mo savings  |  Requires: Redis Enterprise + Embeddings model",
                        font_size=13, bold=True, color=WHITE, bg_color=GREEN,
                        alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Semantic Caching
The MS Learn diagram on the left shows the full semantic caching flow through API Management.
This is the biggest ROI feature. Uses embeddings + vector search to find semantically similar prompts.
DEMO: Run 04-semantic-caching-demo.py to show identical and similar questions returning cached results.
Cost projection: Patient FAQ with 80% cache hit saves ~$360/month.
Requirements: BasicV2+ tier, Redis Enterprise with RediSearch, embeddings model (text-embedding-3-small).
Limitation: Does NOT work with streaming responses.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 10 — Content Safety (MS Learn image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "AI Policy: Content Safety & Security", RED)

    # MS Learn content safety diagram
    try:
        slide.shapes.add_picture(img("content-safety.png"),
                                  Inches(0.5), Inches(1.2), Inches(7.0))
    except Exception:
        pass
    add_text_box(slide, Inches(0.5), Inches(5.8), Inches(7.0), Inches(0.3),
                 "Source: Microsoft Learn — Content safety policy in API Management",
                 font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # RIGHT — security features
    add_text_box(slide, Inches(8.0), Inches(1.2), Inches(4.8), Inches(0.4),
                 "Security Layers", font_size=18, color=RED, bold=True)

    security_cards = [
        ("🔐  Managed Identity", "APIM authenticates to Azure OpenAI\nusing its system-assigned identity.\nNo API keys stored anywhere."),
        ("🛡️  Content Safety", "Server-side filtering: hate, violence,\njailbreak, self-harm. Blocks BEFORE\nthe request reaches the model."),
        ("🔑  OAuth / Credential Mgr", "Authorize AI apps and agents\nusing credential manager.\nSecure access to MCP servers."),
    ]
    for i, (title, desc) in enumerate(security_cards):
        y = Inches(1.8 + i * 1.45)
        add_card(slide, Inches(8.0), y, Inches(4.8), Inches(1.3),
                 title, desc.split("\n"),
                 title_bg=RED, body_bg=RGBColor(0xFD, 0xED, 0xED))

    set_notes(slide, """SPEAKER NOTES — Content Safety
The MS Learn diagram shows how the content safety policy works in API Management.
Three security layers:
1. Managed Identity — APIM authenticates to Azure OpenAI using managed identity. No API keys stored in any app.
2. Content Safety — Azure AI Content Safety service filters prompts for hate, violence, jailbreak, self-harm. Blocks before hitting the model.
3. OAuth / Credential Manager — Authorize AI apps using API Management's credential manager for MCP servers and agents.
For healthcare: content safety is essential for compliance. Block inappropriate queries before they reach GPT.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 11 — Resiliency: Load Balancing + Circuit Breaker (MS Learn images)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Resiliency: Load Balancing & Circuit Breaker", TEAL)

    # LEFT — Load Balancer image
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.4),
                 "Backend Load Balancer", font_size=18, color=TEAL, bold=True)
    try:
        slide.shapes.add_picture(img("backend-load-balancing.png"),
                                  Inches(0.5), Inches(1.8), Inches(5.8))
    except Exception:
        pass
    lb_items = [
        "• Round-robin, weighted, priority-based",
        "• Session-aware distribution",
        "• Optimal PTU utilization with priorities",
    ]
    txBox = add_text_box(slide, Inches(0.5), Inches(5.3), Inches(5.8), Inches(1.5), "")
    add_bullet_list(txBox.text_frame, lb_items, font_size=13, color=DARK_GRAY)

    # RIGHT — Circuit Breaker image
    add_text_box(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.4),
                 "Circuit Breaker", font_size=18, color=TEAL, bold=True)
    try:
        slide.shapes.add_picture(img("backend-circuit-breaker.png"),
                                  Inches(6.8), Inches(1.8), Inches(5.8))
    except Exception:
        pass
    cb_items = [
        "• Dynamic trip duration from Retry-After header",
        "• Automatic failover to healthy backends",
        "• Maximizes priority backend utilization",
    ]
    txBox = add_text_box(slide, Inches(6.8), Inches(5.3), Inches(5.8), Inches(1.5), "")
    add_bullet_list(txBox.text_frame, cb_items, font_size=13, color=DARK_GRAY)

    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.3),
                 "Source: Microsoft Learn — API Management backends",
                 font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Resiliency
Two key resiliency features shown with official Microsoft Learn diagrams:
1. Backend Load Balancer — supports round-robin, weighted, priority-based, and session-aware. Critical for PTU optimization.
2. Circuit Breaker — uses the Retry-After header from throttled backends for dynamic trip duration. When one OpenAI endpoint returns 429, traffic routes to healthy backends automatically.
This is crucial at enterprise scale: multiple AI endpoints, some PTU (reserved), some PAYGO. The gateway optimizes utilization.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 12 — Observability & Token Metrics (MS Learn images)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Observability: Token Metrics & Dashboards", ORANGE)

    # LEFT — emit token metrics diagram
    try:
        slide.shapes.add_picture(img("emit-token-metrics.png"),
                                  Inches(0.3), Inches(1.2), Inches(6.0))
    except Exception:
        pass
    add_text_box(slide, Inches(0.3), Inches(5.0), Inches(6.0), Inches(0.3),
                 "Emitting token metrics through APIM", font_size=10, color=MED_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # RIGHT — analytics workbook screenshot
    try:
        slide.shapes.add_picture(img("analytics-workbook.png"),
                                  Inches(6.8), Inches(1.2), Inches(6.0))
    except Exception:
        pass
    add_text_box(slide, Inches(6.8), Inches(5.0), Inches(6.0), Inches(0.3),
                 "Built-in LLM analytics dashboard in Azure Portal", font_size=10, color=MED_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # Policy snippet
    metric_xml = ('<llm-emit-token-metric namespace="llm-metrics">\n'
                  '    <dimension name="Client IP" value="@(context.Request.IpAddress)" />\n'
                  '    <dimension name="API ID" value="@(context.Api.Id)" />\n'
                  '    <dimension name="User ID" value="@(context.Request.Headers.GetValueOrDefault(\n'
                  '        \\"x-user-id\\", \\"N/A\\"))" />\n'
                  '</llm-emit-token-metric>')
    add_shape_with_text(slide, Inches(0.3), Inches(5.5), Inches(12.5), Inches(1.5),
                        metric_xml, font_size=12, color=LIGHT_BLUE, bg_color=DARK_BG,
                        alignment=PP_ALIGN.LEFT)

    set_notes(slide, """SPEAKER NOTES — Observability
Two MS Learn visuals here:
LEFT: Shows how the llm-emit-token-metric policy emits token consumption data with custom dimensions.
RIGHT: The built-in analytics dashboard in Azure Portal showing token usage patterns across all AI APIs.
The policy XML shown at bottom adds three dimensions: Client IP, API ID, and a custom User ID header.
This is how you answer: 'Which team consumed the most tokens this week? Which model?' — all from Azure Monitor.
Combined with Application Insights, you get full observability: latency, errors, token usage, prompts/completions.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 13 — Gateway Policy XML (code slide)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 "Gateway Policy — Full XML Configuration", font_size=28, color=WHITE, bold=True)

    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(0.3), Inches(0.08), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    policy_code = """<policies>
  <inbound>
    <base />
    <!-- 1. Route to Azure OpenAI backend -->
    <set-backend-service backend-id="openai-backend" />

    <!-- 2. Authenticate with Managed Identity (no API keys!) -->
    <authentication-managed-identity
        resource="https://cognitiveservices.azure.com" />

    <!-- 3. Token Rate Limiting: 10K TPM per subscription -->
    <azure-openai-token-limit
        tokens-per-minute="10000"
        counter-key="@(context.Subscription.Id)"
        estimate-prompt-tokens="true" />

    <!-- 4. Emit Token Metrics to Azure Monitor -->
    <azure-openai-emit-token-metric namespace="AzureOpenAI">
        <dimension name="Subscription" />
        <dimension name="Model" />
    </azure-openai-emit-token-metric>
  </inbound>
  <outbound>
    <base />
    <set-header name="x-apim-request-id" exists-action="override">
        <value>@(context.RequestId.ToString())</value>
    </set-header>
  </outbound>
</policies>"""

    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(7.5), Inches(6.0),
                 policy_code, font_size=13, color=LIGHT_BLUE)

    # RIGHT — annotated steps
    annotations = [
        (Inches(1.6), "① Route", "set-backend-service", BLUE),
        (Inches(2.5), "② Auth", "managed-identity", GREEN),
        (Inches(3.5), "③ Limit", "token-limit (TPM)", ORANGE),
        (Inches(4.8), "④ Metrics", "emit-token-metric", TEAL),
        (Inches(5.8), "⑤ Trace", "x-apim-request-id", MED_GRAY),
    ]
    for y, step, desc, clr in annotations:
        add_shape_with_text(slide, Inches(8.8), y, Inches(4.0), Inches(0.7),
                            f"{step}\n{desc}", font_size=13, bold=True,
                            color=WHITE, bg_color=clr, alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Policy XML
Walk through the annotated XML:
1. set-backend-service — routes to the Azure OpenAI backend we configured
2. authentication-managed-identity — APIM authenticates to Azure OpenAI using its managed identity. No API keys!
3. azure-openai-token-limit — 10K tokens per minute per subscription. Pre-calculates prompt tokens.
4. azure-openai-emit-token-metric — sends token counts to Azure Monitor with Subscription and Model dimensions.
5. Outbound: injects x-apim-request-id for end-to-end tracing.
Point out: this is ALL you need. No code changes in the apps — just change the endpoint URL.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 14 — Live Demo: Setup APIM
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Part 3: Hands-On — Setting Up the AI Gateway", TEAL)

    # Step cards in a horizontal flow
    steps = [
        ("①", "Create APIM", "Consumption tier\n$0 base cost\nManaged Identity ON", TEAL),
        ("②", "Grant RBAC", "Cognitive Services\nUser role\nNo API keys", GREEN),
        ("③", "Add Backend", "Azure OpenAI\nendpoint URL\nHTTP protocol", BLUE),
        ("④", "Import API", "OpenAI inference\nAPI spec\n/openai path", ACCENT_BLUE),
        ("⑤", "Apply Policies", "Token limits\nMetrics emission\n→ AI Gateway!", ORANGE),
    ]

    for i, (num, title, desc, clr) in enumerate(steps):
        x = Inches(0.3 + i * 2.6)
        # Number circle
        add_shape_with_text(slide, x + Inches(0.7), Inches(1.3), Inches(0.7), Inches(0.7),
                            num, font_size=22, bold=True, color=WHITE, bg_color=clr,
                            alignment=PP_ALIGN.CENTER, shape_type=MSO_SHAPE.OVAL)
        # Title card
        add_card(slide, x, Inches(2.2), Inches(2.4), Inches(2.8),
                 title, desc.split("\n"),
                 title_bg=clr, body_bg=CARD_BG,
                 title_size=16, body_size=13)
        # Arrow between steps
        if i < len(steps) - 1:
            add_text_box(slide, x + Inches(2.42), Inches(3.2), Inches(0.3), Inches(0.4),
                         "▶", font_size=18, color=clr, bold=True)

    # Resources bar
    add_shape_with_text(slide, Inches(0.3), Inches(5.3), Inches(12.7), Inches(0.6),
                        "Resources: APIM Consumption (~$3.50/M calls)  •  AI Services S0 (pay-per-token)  •  App Insights (free tier)",
                        font_size=13, color=WHITE, bg_color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # Demo callout
    add_shape_with_text(slide, Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.8),
                        "🟢  LIVE DEMO: 02-setup-apim-gateway.py — Watch the gateway come to life in < 5 minutes",
                        font_size=16, bold=True, color=WHITE, bg_color=GREEN,
                        alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Setup APIM
DEMO: Run 02-setup-apim-gateway.py.
Walk through each step as it executes:
Step 1: Consumption tier = $0 base cost. Enable system-assigned managed identity.
Step 2: RBAC grants 'Cognitive Services User' on the AI resource. No API keys.
Step 3: Backend definition tells APIM where Azure OpenAI lives.
Step 4: Import the OpenAI API spec so APIM knows the schema.
Step 5: Apply AI policies — this is the transformation from APIM to AI Gateway.
Emphasize: the whole setup is automated and reproducible.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 15 — The Transformation: Before vs After
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Part 4: The Transformation — Governed AI Calls", GREEN)

    # BEFORE card
    add_card(slide, Inches(0.3), Inches(1.2), Inches(6.0), Inches(4.5),
             "❌  Before — No Gateway",
             ["", "🔑  API keys in every app",
              "💸  Unlimited token burn",
              "🔄  Same question = same cost every time",
              "📊  No metrics, no dashboard",
              "⚠️  No content safety",
              "🔍  No per-app cost attribution",
              "", "Result: Uncontrolled AI spend"],
             title_bg=RED, body_bg=RGBColor(0xFD, 0xED, 0xED),
             title_size=18, body_size=14, body_color=DARK_GRAY)

    # Arrow
    add_text_box(slide, Inches(6.4), Inches(3.0), Inches(0.8), Inches(0.6),
                 "▶▶", font_size=28, color=GREEN, bold=True)

    # AFTER card
    add_card(slide, Inches(7.0), Inches(1.2), Inches(6.0), Inches(4.5),
             "✅  After — AI Gateway Active",
             ["", "🔐  Subscription keys (never AI API keys)",
              "⏱️  10K tokens/min per app",
              "💾  Semantic cache → 60-80% savings",
              "📊  Token metrics per app, per model",
              "🛡️  Content safety filters active",
              "🆔  x-apim-request-id on every call",
              "", "Result: Governed, observable, cost-efficient"],
             title_bg=GREEN, body_bg=RGBColor(0xE8, 0xF5, 0xE9),
             title_size=18, body_size=14, body_color=DARK_GRAY)

    # Demo callout
    add_shape_with_text(slide, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.8),
                        "🟢  LIVE DEMO: 03-with-apim.py — Same 10 questions, now with full governance. Compare the difference!",
                        font_size=16, bold=True, color=WHITE, bg_color=GREEN,
                        alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — The Transformation
DEMO: Run 03-with-apim.py. Same 10 healthcare questions from Part 1.
Point out the APIM headers in every response: x-apim-request-id, x-apim-model, x-apim-region.
Compare side by side: same questions, but now with tracking, rate limiting, and audit trail.
The app code is nearly identical — only changed the endpoint URL and added a subscription key header.
No SDK changes, no code refactoring. That's the power of a gateway pattern.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 16 — Semantic Caching Deep Dive (LIVE DEMO slide)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Live Demo: Semantic Caching in Action", GREEN)

    # Cost comparison table
    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.4),
                 "Cost Savings Projection — Healthcare Workloads", font_size=18, color=GREEN, bold=True)

    # Table
    th_data = ["Workload", "Calls/Day", "Cache Hit", "Monthly Savings"]
    th_x = [Inches(0.5), Inches(3.0), Inches(5.0), Inches(7.0)]
    th_w = [Inches(2.3), Inches(1.8), Inches(1.8), Inches(2.3)]
    for x, w, h in zip(th_x, th_w, th_data):
        add_shape_with_text(slide, x, Inches(1.8), w, Inches(0.45),
                            h, font_size=12, bold=True, color=WHITE, bg_color=GREEN,
                            alignment=PP_ALIGN.CENTER)
    cost_rows = [
        ("Patient FAQ", "5,000", "80%", "$360/mo"),
        ("Clinical Decision", "2,000", "60%", "$108/mo"),
        ("Admin Reports", "1,000", "70%", "$63/mo"),
        ("Mixed Total", "8,000", "72% avg", "$531/mo saved"),
    ]
    for i, (wl, calls, hit, savings) in enumerate(cost_rows):
        y = Inches(2.35 + i * 0.5)
        bg = BG_GRAY if i % 2 == 0 else WHITE
        bld = (i == len(cost_rows) - 1)
        for x, w, val in zip(th_x, th_w, [wl, calls, hit, savings]):
            c = GREEN if val.startswith("$") else DARK_GRAY
            add_shape_with_text(slide, x, y, w, Inches(0.45),
                                val, font_size=12, bold=bld, color=c, bg_color=bg,
                                alignment=PP_ALIGN.CENTER)

    # Right side — cache flow visual
    add_text_box(slide, Inches(9.8), Inches(1.2), Inches(3.2), Inches(0.4),
                 "Cache Flow", font_size=16, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    flow_boxes = [
        ("Prompt Arrives", BLUE),
        ("Generate Embedding", ACCENT_BLUE),
        ("Search Redis Cache", TEAL),
        ("Hit ≥ 80%? → Return", GREEN),
        ("Miss → Call OpenAI", ORANGE),
        ("Store in Cache", GREEN),
    ]
    for i, (label, clr) in enumerate(flow_boxes):
        y = Inches(1.8 + i * 0.75)
        add_shape_with_text(slide, Inches(9.8), y, Inches(3.2), Inches(0.55),
                            label, font_size=12, bold=True, color=WHITE, bg_color=clr,
                            alignment=PP_ALIGN.CENTER)
        if i < len(flow_boxes) - 1:
            add_text_box(slide, Inches(11.1), y + Inches(0.55), Inches(0.5), Inches(0.2),
                         "▼", font_size=12, color=clr, bold=True, alignment=PP_ALIGN.CENTER)

    # Demo callout
    add_shape_with_text(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8),
                        "🟢  LIVE DEMO: 04-semantic-caching-demo.py — Watch identical & similar questions return cached results",
                        font_size=16, bold=True, color=WHITE, bg_color=GREEN,
                        alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Semantic Caching Demo
DEMO: Run 04-semantic-caching-demo.py.
Show two scenarios:
1. Exact same question 5 times — after first call, responses come from cache with near-zero latency.
2. Similar questions — different wording but same intent all hit the cache.
Cost projection table shows real healthcare workloads. Patient FAQ with 80% cache hit saves ~$360/month.
Redis Enterprise costs ~$200/month but saves $531+/month in tokens at scale. Clear ROI.
Requirements: BasicV2+ tier, Redis Enterprise with RediSearch, embeddings model (text-embedding-3-small).
Limitation: Does NOT work with streaming responses.""")
    # ══════════════════════════════════════════════════════════
    # SLIDE 17 — Content Understanding + Patient Discharge
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Part 5: Content Understanding — Patient Discharge", ORANGE)

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
                 "The AI Gateway governs ALL AI services — not just LLMs",
                 font_size=17, color=ORANGE, bold=True)

    # LEFT — What is Content Understanding
    add_card(slide, Inches(0.3), Inches(1.8), Inches(4.0), Inches(3.0),
             "Azure Content Understanding", [
                 "AI service for documents, images,",
                 "audio, video analysis",
                 "",
                 "• Custom analyzer schemas",
                 "• Structured field extraction",
                 "• Deployed in westus",
                 "• Same APIM gateway & keys",
             ], title_bg=ORANGE, body_bg=RGBColor(0xFF, 0xF3, 0xE0))

    # CENTER — Pipeline flow
    add_text_box(slide, Inches(4.6), Inches(1.8), Inches(3.5), Inches(0.4),
                 "Patient Discharge Pipeline", font_size=16, color=ORANGE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    pipeline_steps = [
        ("📄  Free-text Summary", MED_GRAY),
        ("▼", None),
        ("🤖  Content Understanding", ORANGE),
        ("▼", None),
        ("📋  Structured JSON", GREEN),
    ]
    for i, (label, clr) in enumerate(pipeline_steps):
        y = Inches(2.4 + i * 0.55)
        if clr:
            add_shape_with_text(slide, Inches(4.6), y, Inches(3.5), Inches(0.45),
                                label, font_size=12, bold=True, color=WHITE, bg_color=clr,
                                alignment=PP_ALIGN.CENTER)
        else:
            add_text_box(slide, Inches(5.8), y, Inches(1.0), Inches(0.3),
                         label, font_size=14, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    # RIGHT — Extracted fields
    add_card(slide, Inches(8.5), Inches(1.8), Inches(4.5), Inches(3.0),
             "Extracted Fields → EHR Integration", [
                 "• Patient Name, MRN, DOB",
                 "• Primary Diagnosis (ICD code)",
                 "• Medications + Dosages",
                 "• Follow-up Instructions",
                 "• Discharge Date",
                 "",
                 "Output: Structured JSON",
             ], title_bg=GREEN, body_bg=RGBColor(0xE8, 0xF5, 0xE9))

    # Bottom — one gateway message
    add_shape_with_text(slide, Inches(0.3), Inches(5.3), Inches(12.7), Inches(1.5),
                        "One Gateway — Multiple AI Services\n\n"
                        "Azure OpenAI (GPT-4o)  •  Content Understanding  •  Document Intelligence\n"
                        "Speech Services  •  Custom Vision  •  All governed through APIM",
                        font_size=15, color=WHITE, bg_color=ORANGE, alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Content Understanding
This elevates the demo beyond just LLM calls. APIM isn't just for GPT.
We provisioned Content Understanding in westus, added as a second backend in APIM.
Use case: Patient discharge summaries are free text. CU parses them into structured JSON — patient name, MRN, diagnosis, medications, follow-up instructions.
All going through the SAME APIM gateway with the SAME subscription keys, SAME rate limiting, SAME metrics.
One dashboard shows ALL AI consumption: LLMs AND document processing AND speech AND vision.
DEMO: Run 02-patient-discharge-analyzer.py if Content Understanding is provisioned.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 18 — Budget & Cost Controls
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Budget & Cost Controls", ACCENT_BLUE)

    # LEFT — Budget gauge visual
    add_card(slide, Inches(0.3), Inches(1.2), Inches(5.5), Inches(4.5),
             "Resource Group Budget: $25 / month",
             ["", "🟢  25%  ─────  $6.25 spent",
              "🟡  50%  ─────  $12.50 spent",
              "🟠  75%  ─────  $18.75 spent — review",
              "🔴  90%  ─────  $22.50 spent — action!",
              "",
              "Email: jaypadhya@microsoft.com",
              "",
              "Subscription Budget: $100/month (safety net)"],
             title_bg=ACCENT_BLUE, body_bg=CARD_BG,
             title_size=18, body_size=14)

    # RIGHT — Cost breakdown cards
    add_text_box(slide, Inches(6.3), Inches(1.2), Inches(6.5), Inches(0.4),
                 "Resource Cost Breakdown", font_size=18, color=ACCENT_BLUE, bold=True)

    costs = [
        ("APIM Consumption", "~$3.50/million calls", "$0 base cost", BLUE),
        ("AI Services (S0)", "Pay-per-token", "GPT-4o: ~$0.03/1K tokens", GREEN),
        ("Content Understanding", "Pay-per-page", "Document processing", ORANGE),
        ("App Insights", "Free tier", "5GB/month included", TEAL),
        ("Redis Enterprise", "~$200/month", "For semantic caching (optional)", RED),
    ]
    for i, (name, price, note, clr) in enumerate(costs):
        y = Inches(1.8 + i * 0.75)
        add_shape_with_text(slide, Inches(6.3), y, Inches(0.12), Inches(0.6),
                            "", bg_color=clr, shape_type=MSO_SHAPE.RECTANGLE)
        add_text_box(slide, Inches(6.6), y, Inches(3.0), Inches(0.35),
                     name, font_size=14, bold=True, color=DARK_BG)
        add_text_box(slide, Inches(6.6), y + Inches(0.3), Inches(3.0), Inches(0.3),
                     note, font_size=11, color=MED_GRAY)
        add_text_box(slide, Inches(10.0), y + Inches(0.1), Inches(2.8), Inches(0.35),
                     price, font_size=13, bold=True, color=clr, alignment=PP_ALIGN.RIGHT)

    # ROI callout
    add_shape_with_text(slide, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.8),
                        "💡  ROI: Semantic caching saves $531+/mo in tokens — more than the cost of Redis Enterprise ($200/mo). The gateway pays for itself.",
                        font_size=15, bold=True, color=WHITE, bg_color=ACCENT_BLUE,
                        alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Budget & Cost Controls
$25/month budget on the resource group with alerts at 25%, 50%, 75%, 90%.
Costs for this demo are minimal: Consumption tier APIM is $0 base, AI calls are pay-per-token.
The cost breakdown shows each component's pricing model.
KEY MESSAGE: The gateway + caching combination can REDUCE your AI spend by 60-80%. Caching saves more than it costs.
Show: 05-setup-budget-alerts.py briefly to demonstrate automated cost governance.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 19 — CI/CD Pipeline
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "CI/CD — GitHub Actions Pipeline", DARK_BG)

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
                 "github.com/jaypadhya1605/foundry-agent-observatory",
                 font_size=15, color=BLUE, bold=True)

    # Pipeline stages — connected boxes
    stages = [
        ("①  Validate", "Lint Python scripts\nValidate XML policies\nSchemacheck configs", BLUE),
        ("②  Deploy Infra", "APIM + AI Gateway\nContent Understanding\nBudget Alerts", GREEN),
        ("③  Integration Test", "Smoke test OpenAI\nSmoke test CU\nHealth checks", ORANGE),
    ]
    for i, (name, desc, clr) in enumerate(stages):
        x = Inches(0.3 + i * 4.3)
        add_card(slide, x, Inches(1.8), Inches(3.8), Inches(2.5),
                 name, desc.split("\n"),
                 title_bg=clr, body_bg=CARD_BG,
                 title_size=18, body_size=14)
        if i < len(stages) - 1:
            add_text_box(slide, x + Inches(3.85), Inches(2.5), Inches(0.5), Inches(0.5),
                         "▶", font_size=24, color=clr, bold=True)

    # Key features
    features = [
        "🔐  OIDC authentication — no stored secrets in GitHub (Federated Identity Credentials)",
        "🎯  Triggered on push to main (folder-2-APIM/** paths) or manual workflow dispatch",
        "✅  XML policy validation catches gateway config errors before deployment",
        "🧪  Smoke tests verify both OpenAI and Content Understanding endpoints post-deploy",
    ]
    for i, feat in enumerate(features):
        y = Inches(4.6 + i * 0.45)
        add_text_box(slide, Inches(0.5), y, Inches(12.5), Inches(0.4),
                     feat, font_size=13, color=DARK_GRAY)

    set_notes(slide, """SPEAKER NOTES — CI/CD
GitHub Actions pipeline with three stages:
1. Validate — lints Python files and validates XML gateway policies. Catches errors before deployment.
2. Deploy Infrastructure — runs APIM setup, CU provisioning, budget alerts. Uses OIDC auth — no secrets.
3. Integration Test — smoke tests both endpoints through APIM.
Every policy change goes through validation before production.
Repo: github.com/jaypadhya1605/foundry-agent-observatory""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 20 — APIM Tier Comparison
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "APIM Tier Comparison — Which One to Choose?", BLUE)

    # Table header
    cols_x = [Inches(0.5), Inches(3.0), Inches(5.0), Inches(7.0), Inches(9.5)]
    cols_w = [Inches(2.3), Inches(1.8), Inches(1.8), Inches(2.3), Inches(2.8)]
    headers = ["Tier", "Base Cost", "Per-Call", "Best For", "AI Features"]
    for x, w, h in zip(cols_x, cols_w, headers):
        add_shape_with_text(slide, x, Inches(1.2), w, Inches(0.5),
                            h, font_size=13, bold=True, color=WHITE, bg_color=DARK_BG,
                            alignment=PP_ALIGN.CENTER)

    tiers = [
        ("Consumption", "$0/mo", "~$3.50/M", "Demos, pilots", "✅ This demo"),
        ("Basic v2", "~$170/mo", "Included", "Small production", "Token policies"),
        ("Standard v2", "~$700/mo", "Included", "Medium prod.", "+ VNet, caching"),
        ("Premium", "~$2,800/mo", "Included", "Enterprise", "+ Multi-region"),
    ]
    for i, row in enumerate(tiers):
        y = Inches(1.85 + i * 0.6)
        bg = BG_GRAY if i % 2 == 0 else WHITE
        for j, (x, w, val) in enumerate(zip(cols_x, cols_w, row)):
            bld = (j == 4 and "This demo" in val) or j == 0
            clr = GREEN if "This demo" in val else DARK_GRAY
            add_shape_with_text(slide, x, y, w, Inches(0.5),
                                val, font_size=12, bold=bld, color=clr, bg_color=bg,
                                alignment=PP_ALIGN.CENTER)

    # Recommendation card
    add_shape_with_text(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                        "Recommendation for Providence\n\n"
                        "①  Start with Consumption tier for this pilot (today's demo — $0 base cost)\n"
                        "②  Enable token metrics + rate limiting immediately (value on any tier)\n"
                        "③  Test semantic caching on Patient FAQ workloads during pilot\n"
                        "④  Graduate to Standard v2 for production (VNet integration, full AI policies, caching)\n"
                        "⑤  Expand to Content Understanding for discharge docs + other AI services",
                        font_size=14, color=WHITE, bg_color=BLUE)

    set_notes(slide, """SPEAKER NOTES — Tier Comparison
Consumption tier is what we're using today. $0 base, pay per call.
For production AI Gateway with full features (token limiting, semantic caching, metrics), need BasicV2+.
Standard v2 adds VNet integration — important for healthcare network security.
Recommendation: Start Consumption to prove concept. Then Standard v2 for production.
Semantic caching alone saves more than the tier upgrade cost.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 21 — Pros & Cons
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_section_divider(slide, Inches(0.8), Inches(0.4), Inches(10),
                        "Honest Assessment — Pros & Considerations", DARK_BG)

    # Pros card
    add_card(slide, Inches(0.3), Inches(1.2), Inches(6.0), Inches(5.2),
             "✅  Advantages",
             ["",
              "💰  Token budgets prevent runaway AI spend",
              "💾  Semantic caching → 60-80% cost savings",
              "🔐  Managed identity eliminates API key sprawl",
              "📊  Centralized metrics across ALL AI services",
              "👥  Multi-tenancy: isolate usage per team/app",
              "🛡️  Content safety filtering layer",
              "🔧  No code changes — just change endpoint URL",
              "",
              "The gateway pays for itself through",
              "caching savings alone."],
             title_bg=GREEN, body_bg=RGBColor(0xE8, 0xF5, 0xE9),
             title_size=18, body_size=14)

    # Cons card
    add_card(slide, Inches(7.0), Inches(1.2), Inches(6.0), Inches(5.2),
             "⚠️  Considerations",
             ["",
              "⏱️  ~10-50ms added latency (negligible)",
              "📡  Streaming: some policies + SSE = complex",
              "💵  Redis Enterprise: ~$200/mo for caching",
              "📖  XML policies: APIM learning curve",
              "🌐  Consumption tier: single-region only",
              "📦  Batch/fine-tuning don't benefit from gateway",
              "",
              "For Providence: high-repetition workloads",
              "(patient FAQ, clinical lookups) are ideal",
              "for semantic caching ROI."],
             title_bg=ORANGE, body_bg=RGBColor(0xFF, 0xF3, 0xE0),
             title_size=18, body_size=14)

    set_notes(slide, """SPEAKER NOTES — Pros & Cons
Be honest about tradeoffs.
PROS: Cost + security story is compelling. Managed identity + token budgets + caching.
CONS: ~10-50ms latency is negligible. Streaming is the main limitation. Redis costs $200/mo but saves $500+/mo.
For Providence: high-repetition workloads (patient FAQ, clinical lookups) are PERFECT for semantic caching.""")

    # ══════════════════════════════════════════════════════════
    # SLIDE 22 — Summary & Next Steps
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(0.5), Inches(0.08), Inches(6.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "Summary & Next Steps", font_size=36, color=WHITE, bold=True)

    # What we demonstrated — left
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
                 "What We Demonstrated Today", font_size=20, color=LIGHT_BLUE, bold=True)
    demos = [
        "✅  Direct vs. governed AI calls (before/after)",
        "✅  AI Gateway capabilities (official MS Learn diagrams)",
        "✅  Token rate limiting, semantic caching, content safety",
        "✅  Load balancing & circuit breaker for resiliency",
        "✅  Observability — token metrics & dashboards",
        "✅  Content Understanding for patient discharge docs",
        "✅  Budget alerts ($25/mo) with automated thresholds",
        "✅  CI/CD pipeline on GitHub Actions",
    ]
    txBox = add_text_box(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.8), "")
    add_bullet_list(txBox.text_frame, demos, font_size=14, color=LIGHT_GRAY)

    # Next steps — right
    add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4),
                 "Recommended Next Steps", font_size=20, color=LIGHT_BLUE, bold=True)
    nexts = [
        "①  Pilot Consumption tier with one workload",
        "②  Enable token metrics (immediate observability)",
        "③  Test semantic caching on Patient FAQ",
        "④  Plan Standard v2 upgrade for production",
        "⑤  Onboard Content Understanding for discharge docs",
        "⑥  Expand to Speech, Vision, Document Intelligence",
        "⑦  Explore AI Gateway in Microsoft Foundry (preview)",
    ]
    txBox = add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(3.8), "")
    add_bullet_list(txBox.text_frame, nexts, font_size=14, color=LIGHT_GRAY)

    # Resources bar
    add_shape_with_text(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9),
                        "Resources\n"
                        "GitHub: github.com/jaypadhya1605/foundry-agent-observatory  |  "
                        "MS Learn: learn.microsoft.com/azure/api-management/genai-gateway-capabilities  |  "
                        "RG: rg-apim-providence-demo-jp-001",
                        font_size=12, color=LIGHT_GRAY, bg_color=RGBColor(0x22, 0x22, 0x22),
                        alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — Summary
Recap:
1. Started with zero governance — apps calling OpenAI directly with API keys.
2. Learned AI Gateway capabilities from official Microsoft Learn documentation.
3. Built an AI Gateway with APIM — automated setup, managed identity, token policies, metrics.
4. Showed semantic caching for 60-80% savings, content safety, and load balancing.
5. Extended to Content Understanding for patient discharge docs.
6. Set up $25/month budget alerts and CI/CD pipeline.
Next steps: Pilot, enable metrics, test caching, graduate to Standard v2.
Thank you!""")

    # ══════════════════════════════════════════════════════════
    # Save
    # ══════════════════════════════════════════════════════════
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "Providence-APIM-AI-Gateway-Deep-Dive.pptx")
    prs.save(output_path)
    print(f"\n✅ Presentation saved: {output_path}")
    print(f"   Slides: {len(prs.slides)}")
    print(f"   All slides have speaker notes.")
    return output_path


if __name__ == "__main__":
    create_presentation()
